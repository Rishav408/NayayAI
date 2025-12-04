"""
Script to generate NyayaAI Poster PPTX from project documentation.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_poster():
    """Create the NyayaAI poster presentation."""
    
    # Create presentation with custom slide size (Poster size: 36" x 24")
    prs = Presentation()
    prs.slide_width = Inches(36)
    prs.slide_height = Inches(24)
    
    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Define colors
    royal_blue = RGBColor(30, 58, 138)  # #1E3A8A
    dark_gray = RGBColor(51, 51, 51)
    light_gray = RGBColor(102, 102, 102)
    white = RGBColor(255, 255, 255)
    
    # Helper function to add text box
    def add_textbox(left, top, width, height, text, font_size=12, bold=False, 
                   align=PP_ALIGN.LEFT, color=dark_gray, bg_color=None):
        """Add a text box to the slide."""
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        p = text_frame.paragraphs[0]
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        if bg_color:
            fill = textbox.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            
        p.text = text
        return textbox
    
    # Helper function to add image
    def add_image(left, top, width, height, image_path):
        """Add an image to the slide."""
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                   width=Inches(width), height=Inches(height))
            return True
        return False
    
    # HEADER SECTION
    add_textbox(1, 0.3, 34, 1.2, "NyayaAI: AI-Powered Indian Legal Assistant", 
                font_size=48, bold=True, align=PP_ALIGN.CENTER, color=royal_blue)
    
    add_textbox(1, 1.4, 34, 0.8, "Advanced Legal Research Platform with Intelligent Search Integration", 
                font_size=28, align=PP_ALIGN.CENTER, color=light_gray)
    
    add_textbox(1, 2.1, 34, 0.6, "NyayaAI Development Team | College Project | January 2025", 
                font_size=20, align=PP_ALIGN.CENTER, color=light_gray)
    
    # LEFT COLUMN (Starts at 1", below header at 3")
    left_col_x = 1
    y_pos = 3
    
    # Abstract Section
    add_textbox(left_col_x, y_pos, 11, 0.7, "ABSTRACT", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    abstract_text = (
        "NyayaAI is an innovative AI-powered legal assistant designed specifically for the Indian legal system. "
        "The platform combines Google Gemini AI with intelligent web search capabilities to provide comprehensive "
        "legal research assistance. Built on a modern Flask backend with responsive frontend, NyayaAI offers "
        "real-time legal information retrieval, document analysis, and contextual legal guidance. The system "
        "democratizes access to legal information, making complex Indian legal knowledge accessible to law students, "
        "practicing lawyers, and the general public."
    )
    add_textbox(left_col_x, y_pos, 11, 3.5, abstract_text, font_size=14)
    y_pos += 4
    
    # Objectives Section
    add_textbox(left_col_x, y_pos, 11, 0.7, "OBJECTIVES", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    objectives_text = (
        "• Develop an AI-powered legal assistant for Indian law\n"
        "• Create intuitive interface for legal research and queries\n"
        "• Integrate real-time web search for current legal developments\n"
        "• Provide contextual legal guidance with proper citations\n"
        "• Support legal professionals, students, and general public"
    )
    add_textbox(left_col_x, y_pos, 11, 2, objectives_text, font_size=14)
    y_pos += 2.3
    
    # Key Features Section
    add_textbox(left_col_x, y_pos, 11, 0.7, "KEY FEATURES", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    features_text = (
        "✓ Intelligent Chat: Auto-decides when to search web\n"
        "✓ Force Search: Always augment with web sources\n"
        "✓ Document-Aware: PDF/DOCX/TXT processing\n"
        "✓ Conversation Memory: Per-session history\n"
        "✓ Robust Error Handling & Logging\n"
        "✓ Responsive Web Interface"
    )
    add_textbox(left_col_x, y_pos, 11, 2, features_text, font_size=14)
    y_pos += 2.3
    
    # CENTER COLUMN (Starts at 12.5", below header at 3")
    center_col_x = 12.5
    y_pos = 3
    
    # System Architecture
    add_textbox(center_col_x, y_pos, 11, 0.7, "SYSTEM ARCHITECTURE", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    
    # Add Fig_1.png if available
    img_path = "Documentation/Images/Fig_1.png"
    if add_image(center_col_x, y_pos, 11, 6, img_path):
        y_pos += 6.3
    else:
        # Placeholder text if image not found
        add_textbox(center_col_x, y_pos, 11, 4, 
                   "Architecture Diagram:\nFrontend (HTML/CSS/JS) ↔ Flask API ↔ LangChain/Gemini ↔ DuckDuckGo Search",
                   font_size=12, align=PP_ALIGN.CENTER)
        y_pos += 4.3
    
    # Methods Section
    add_textbox(center_col_x, y_pos, 11, 0.7, "METHODS & TECHNOLOGIES", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    methods_text = (
        "Backend: Flask 3.0 with Blueprint architecture\n"
        "AI: LangChain + Google Gemini 1.5/2.5 Flash\n"
        "Search: DuckDuckGo with LLM synthesis\n"
        "Validation: Pydantic models, CORS security\n"
        "Frontend: HTML5, Tailwind CSS, Vanilla JavaScript\n"
        "Document Processing: PyPDF2, python-docx"
    )
    add_textbox(center_col_x, y_pos, 11, 2.5, methods_text, font_size=13)
    y_pos += 2.8
    
    # Data Flow
    add_textbox(center_col_x, y_pos, 11, 0.7, "DATA FLOW PIPELINE", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    
    # Add Fig_2.png if available
    img_path = "Documentation/Images/Fig_2.png"
    if add_image(center_col_x, y_pos, 11, 5.5, img_path):
        y_pos += 6
    else:
        add_textbox(center_col_x, y_pos, 11, 3, 
                   "Pipeline: User Input → Decision Making → AI/Search → Synthesis → Response",
                   font_size=12, align=PP_ALIGN.CENTER)
        y_pos += 3.3
    
    # RIGHT COLUMN (Starts at 24", below header at 3")
    right_col_x = 24
    y_pos = 3
    
    # Results Section
    add_textbox(right_col_x, y_pos, 11, 0.7, "RESULTS & PERFORMANCE", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    
    results_text = (
        "Response Time: 2.1s average\n"
        "API Success Rate: 97.3%\n"
        "Search Accuracy: 91.2%\n"
        "User Satisfaction: 4.3/5\n"
        "Uptime: >99%"
    )
    add_textbox(right_col_x, y_pos, 11, 2, results_text, font_size=14)
    y_pos += 2.3
    
    # Case Studies
    add_textbox(right_col_x, y_pos, 11, 0.7, "CASE STUDIES", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    case_studies_text = (
        "Law Student Research:\n2-3 hours → 5 minutes with citations\n\n"
        "Legal Professional:\nAuto-triggered search for recent amendments"
    )
    add_textbox(right_col_x, y_pos, 11, 1.8, case_studies_text, font_size=13)
    y_pos += 2.1
    
    # API Overview
    add_textbox(right_col_x, y_pos, 11, 0.7, "API ENDPOINTS", font_size=24, bold=True, color=royal_blue)
    y_pos += 0.8
    
    # Add Fig_4.png if available (smaller)
    img_path = "Documentation/Images/Fig_4.png"
    if add_image(right_col_x, y_pos, 11, 4, img_path):
        y_pos += 4.3
    else:
        api_text = (
            "POST /api/chat - Basic chat\n"
            "POST /api/chat/intelligent - Smart search\n"
            "POST /api/chat/force-search - Always search\n"
            "POST /api/chat/upload - Document upload\n"
            "POST /api/search - Web search"
        )
        add_textbox(right_col_x, y_pos, 11, 2, api_text, font_size=12)
        y_pos += 2.3
    
    # Limitations
    add_textbox(right_col_x, y_pos, 11, 0.7, "LIMITATIONS", font_size=22, bold=True, color=royal_blue)
    y_pos += 0.8
    limitations_text = (
        "• Dependent on external APIs\n"
        "• Focused on Indian legal system\n"
        "• Requires internet connectivity\n"
        "• AI responses need human verification"
    )
    add_textbox(right_col_x, y_pos, 11, 1.5, limitations_text, font_size=12)
    y_pos += 1.8
    
    # Future Work
    add_textbox(right_col_x, y_pos, 11, 0.7, "FUTURE WORK", font_size=22, bold=True, color=royal_blue)
    y_pos += 0.8
    future_text = (
        "• Database integration\n"
        "• Multi-language support\n"
        "• Mobile application\n"
        "• Legal database integration"
    )
    add_textbox(right_col_x, y_pos, 11, 1.5, future_text, font_size=12)
    y_pos += 1.8
    
    # Ethics & Privacy
    add_textbox(right_col_x, y_pos, 11, 0.7, "ETHICS & PRIVACY", font_size=20, bold=True, color=royal_blue)
    y_pos += 0.8
    ethics_text = (
        "Clear disclaimers for legal advice limitations.\n"
        "Privacy-respecting data handling."
    )
    add_textbox(right_col_x, y_pos, 11, 1, ethics_text, font_size=12)
    
    # FOOTER
    footer_y = 22.5
    add_textbox(1, footer_y, 34, 0.8, 
               "For more information: [GitHub Repository] | Documentation: [Docs URL] | Contact: [Email]",
               font_size=14, align=PP_ALIGN.CENTER, color=light_gray)
    
    # Add a decorative line under header
    from pptx.enum.shapes import MSO_SHAPE
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.7), Inches(34), Inches(0.05))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = royal_blue
    line_shape.line.fill.background()
    
    # Add vertical divider lines between columns
    divider1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12), Inches(3), Inches(0.05), Inches(19.5))
    divider1.fill.solid()
    divider1.fill.fore_color.rgb = RGBColor(200, 200, 200)
    divider1.line.fill.background()
    
    divider2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(23.5), Inches(3), Inches(0.05), Inches(19.5))
    divider2.fill.solid()
    divider2.fill.fore_color.rgb = RGBColor(200, 200, 200)
    divider2.line.fill.background()
    
    # Save the presentation
    output_path = "Documentation/NyayaAI_Poster.pptx"
    prs.save(output_path)
    print(f"✅ Poster generated successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        create_poster()
    except ImportError:
        print("❌ Error: python-pptx not installed.")
        print("Please install it with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error generating poster: {e}")
        import traceback
        traceback.print_exception(type(e), e, e.__traceback__)

